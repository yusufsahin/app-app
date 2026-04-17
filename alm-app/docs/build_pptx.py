"""
Build Pamera executive pitch deck PPTX.
Provera branding, corporate tone, no startup feel.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ──────────────────────────────────────────────────────────────────
DOCS = Path(__file__).parent
SS   = DOCS / "screenshots"
OUT  = DOCS / "pamera-pitch-deck.pptx"

# ── Brand colours ──────────────────────────────────────────────────────────
BLUE      = RGBColor(0x00, 0x52, 0xCC)
BLUE_DK   = RGBColor(0x00, 0x3D, 0x99)
BLUE_LT   = RGBColor(0xE6, 0xEF, 0xFF)
INK       = RGBColor(0x11, 0x18, 0x27)
BODY      = RGBColor(0x37, 0x41, 0x51)
MUTED     = RGBColor(0x6B, 0x72, 0x80)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BG_ALT    = RGBColor(0xF4, 0xF7, 0xFB)
RED_LT    = RGBColor(0xB9, 0x1C, 0x1C)
GREEN     = RGBColor(0x15, 0x80, 0x3D)
BORDER    = RGBColor(0xD1, 0xD9, 0xE0)

# ── Slide dimensions (Widescreen 16:9) ────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(11), bold=False, color=INK,
             align=PP_ALIGN.LEFT, font_name="Calibri",
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.name      = font_name
    run.font.italic    = italic
    return txb

def add_image(slide, path, l, t, w, h=None):
    if h:
        slide.shapes.add_picture(str(path), l, t, w, h)
    else:
        slide.shapes.add_picture(str(path), l, t, w)

def slide_stripe(slide, color=BLUE, top=True):
    """5px top or bottom accent stripe."""
    y = Inches(0) if top else H - Inches(0.07)
    add_rect(slide, Inches(0), y, W, Inches(0.07), fill=color)

def slide_footer(slide, page_num, total, label=""):
    add_rect(slide, Inches(0), H - Inches(0.45), W, Inches(0.45), fill=BG_ALT)
    add_text(slide, "Provera · Pamera · info@provera.com.tr · provera.com.tr",
             Inches(0.4), H - Inches(0.38), Inches(8), Inches(0.32),
             font_size=Pt(7), color=MUTED, italic=True)
    add_text(slide, f"{page_num} / {total}",
             Inches(12.5), H - Inches(0.38), Inches(0.8), Inches(0.32),
             font_size=Pt(7), color=MUTED, align=PP_ALIGN.RIGHT)

def section_title(slide, text, l=Inches(0.5), t=Inches(0.9),
                  w=Inches(12.3), color=INK):
    add_text(slide, text, l, t, w, Inches(0.45),
             font_size=Pt(18), bold=True, color=color)
    # underline bar
    add_rect(slide, l, t + Inches(0.44), Inches(1.0), Inches(0.03), fill=BLUE)

def bullet_block(slide, items, l, t, w, h_per=Inches(0.35),
                 color=BODY, size=Pt(10.5), bullet="›"):
    for i, item in enumerate(items):
        add_text(slide, f"{bullet}  {item}",
                 l, t + i * h_per, w, h_per,
                 font_size=size, color=color)

def risk_card(slide, l, t, w, h, title, body):
    add_rect(slide, l, t, Inches(0.05), h, fill=RED_LT)
    add_rect(slide, l + Inches(0.05), t, w - Inches(0.05), h,
             fill=RGBColor(0xFE, 0xF2, 0xF2))
    add_text(slide, title, l + Inches(0.15), t + Inches(0.08),
             w - Inches(0.2), Inches(0.28), font_size=Pt(10), bold=True, color=RED_LT)
    add_text(slide, body, l + Inches(0.15), t + Inches(0.32),
             w - Inches(0.2), h - Inches(0.38), font_size=Pt(9), color=BODY)

def value_card(slide, l, t, w, h, title, body):
    add_rect(slide, l, t, w, h,
             fill=RGBColor(0xFF,0xFF,0xFF), line=BORDER, line_w=Pt(0.75))
    add_rect(slide, l, t, Inches(0.05), h, fill=BLUE)
    add_text(slide, title, l + Inches(0.15), t + Inches(0.08),
             w - Inches(0.2), Inches(0.28), font_size=Pt(10), bold=True, color=BLUE_DK)
    add_text(slide, body, l + Inches(0.15), t + Inches(0.32),
             w - Inches(0.2), h - Inches(0.38), font_size=Pt(9), color=BODY)

def outcome_card(slide, l, t, w, h, title, today, gain, outcome):
    add_rect(slide, l, t, w, h,
             fill=WHITE, line=BORDER, line_w=Pt(0.75))
    add_rect(slide, l, t, Inches(0.05), h, fill=BLUE)
    cy = t + Inches(0.1)
    add_text(slide, title, l+Inches(0.15), cy, w-Inches(0.2), Inches(0.28),
             font_size=Pt(10), bold=True, color=INK)
    cy += Inches(0.32)
    add_text(slide, f"Bugün: {today}", l+Inches(0.15), cy, w-Inches(0.2), Inches(0.35),
             font_size=Pt(8.5), color=RED_LT, italic=True)
    cy += Inches(0.38)
    add_text(slide, gain, l+Inches(0.15), cy, w-Inches(0.2), h-Inches(0.9),
             font_size=Pt(9), color=BODY)
    add_text(slide, f"→  {outcome}",
             l+Inches(0.15), t+h-Inches(0.3), w-Inches(0.2), Inches(0.28),
             font_size=Pt(8.5), bold=True, color=GREEN)

TOTAL = 9

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# Dark left panel
add_rect(s, Inches(0), Inches(0), Inches(5), H, fill=INK)
# Blue accent bar inside panel
add_rect(s, Inches(0), Inches(0), Inches(0.08), H, fill=BLUE)

# Logo placeholder
try:
    logo_path = DOCS / "provera_logo.png"
    if logo_path.exists():
        add_image(s, logo_path, Inches(0.5), Inches(0.5), Inches(2.2))
except:
    add_text(s, "PROVERA", Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
             font_size=Pt(16), bold=True, color=WHITE)

# Product name
add_text(s, "Pamera",
         Inches(0.5), Inches(2.2), Inches(4.2), Inches(1.0),
         font_size=Pt(48), bold=True, color=WHITE, font_name="Calibri")

# Tagline
add_text(s, "Yazılım Geliştirme Yönetim Platformu",
         Inches(0.5), Inches(3.2), Inches(4.2), Inches(0.5),
         font_size=Pt(13), color=RGBColor(0xA0,0xB4,0xD0))

# Divider line
add_rect(s, Inches(0.5), Inches(3.85), Inches(1.5), Inches(0.03), fill=BLUE)

# Subtitle
add_text(s, "Planlama · Kalite · Teslimat İzlenebilirliği",
         Inches(0.5), Inches(4.05), Inches(4.2), Inches(0.4),
         font_size=Pt(10), color=RGBColor(0x8A,0x9B,0xB8))

# Company / contact bottom left
add_text(s, "Provera Telekomünikasyon ve Bilgi Teknolojileri",
         Inches(0.5), H - Inches(1.1), Inches(4.2), Inches(0.35),
         font_size=Pt(8.5), color=RGBColor(0x7A,0x8A,0xA0))
add_text(s, "info@provera.com.tr  ·  provera.com.tr",
         Inches(0.5), H - Inches(0.75), Inches(4.2), Inches(0.35),
         font_size=Pt(8.5), color=RGBColor(0x6A,0x7A,0x90))
add_text(s, "Gizli — Yalnızca Alıcı İçin",
         Inches(0.5), H - Inches(0.42), Inches(4.2), Inches(0.28),
         font_size=Pt(7.5), color=RGBColor(0x50,0x60,0x75), italic=True)

# Right panel – screenshot of projects
try:
    add_image(s, SS / "01_projects.png",
              Inches(5.1), Inches(0.3), Inches(8.0), Inches(6.9))
except:
    pass

# subtle overlay on screenshot
add_rect(s, Inches(5.1), Inches(0.3), Inches(8.0), Inches(6.9),
         fill=None)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_stripe(s)
section_title(s, "Neden Önemli — Üç Temel İş Riski")
add_text(s, "Yazılım geliştiren organizasyonların büyük çoğunluğunda planlama, kalite güvencesi ve teslimat takibi birbirinden "
            "kopuk araçlarda yürütülüyor. Bu parçalanma üç somut iş riskine dönüşüyor.",
         Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.55),
         font_size=Pt(10.5), color=BODY)

cw = Inches(3.9)
ch = Inches(3.5)
ct = Inches(2.1)
gap = Inches(0.25)

risk_card(s, Inches(0.5), ct, cw, ch,
          "Risk 1 — Görünürlük Kaybı",
          "Sprint ve kalite durumunu tek yerden görmek için araçlar arası veri toplama saatler alıyor. "
          "Yönetim soruları hızla yanıtlanamıyor; karar döngüsü uzuyor.")

risk_card(s, Inches(0.5) + cw + gap, ct, cw, ch,
          "Risk 2 — Gizli Maliyet Birikimi",
          "Çoklu araç lisansları, entegrasyon bakımı ve araçlar arası geçiş zaman kaybı "
          "bütçede tek satırda görünmüyor — ama her ay biriyor. "
          "30–50 kişilik ekiplerde bu yük yıllık 16.000–24.000 USD'ye ulaşabiliyor.")

risk_card(s, Inches(0.5) + 2*(cw + gap), ct, cw, ch,
          "Risk 3 — Denetim & Uyum Açığı",
          "Gereksinim–test–deployment zincirini el emeğiyle kurmak, her denetim hazırlığını "
          "ayrı bir projeye dönüştürüyor. ISO, ASPICE veya müşteri denetimleri "
          "için harcanan süre tekrarlanıyor.")

slide_footer(s, 2, TOTAL)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_stripe(s)
section_title(s, "Çözüm — Üç Süreç, Tek Platform")

add_text(s,
         "Pamera; planlama, kalite operasyonu ve teslimat izlenebilirliğini tek veri modeli üzerinde birleştirir. "
         "Araçlar arası entegrasyon maliyeti, veri kopukluğu ve görünürlük kaybı ortadan kalkar.",
         Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.55),
         font_size=Pt(10.5), color=BODY)

vw = Inches(2.9)
vh = Inches(3.6)
vt = Inches(2.15)
vgap = Inches(0.25)

value_card(s, Inches(0.5),              vt, vw, vh,
           "Yönetim Görünürlüğü",
           "Sprint durumu, kalite metrikleri ve teslimat geçmişi tek ekrandan anlık okunur. "
           "Yönetim katmanı süreci değil sonuçları izler.")

value_card(s, Inches(0.5)+vw+vgap,     vt, vw, vh,
           "Kalite Operasyonu — Çekirdekte",
           "Test, hata yönetimi ve backlog aynı veri modelinde çalışır. "
           "Ayrı araç lisansı ve entegrasyon bakımı gerekmez.")

value_card(s, Inches(0.5)+2*(vw+vgap), vt, vw, vh,
           "Uçtan Uca İzlenebilirlik",
           "Kod değişiklikleri ve deployment kayıtları ilgili gereksinime otomatik bağlanır. "
           "Denetim hazırlığı için el emeği gerekmez.")

value_card(s, Inches(0.5)+3*(vw+vgap), vt, vw, vh,
           "Sürece Uyan Araç",
           "Her projenin iş akışı yapılandırma dosyası ile tanımlanır. "
           "Araç sürece uyar; ekip araca değil.")

slide_footer(s, 3, TOTAL)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BACKLOG & BOARD
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_stripe(s)
section_title(s, "Backlog & Kanban Board")
add_text(s, "İş öğesi yönetimi ve operasyonel çalışma alanı — tek ekrandan proje görünürlüğü",
         Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.35),
         font_size=Pt(10.5), color=BODY)

# Two screenshots side by side
try:
    add_image(s, SS/"03_backlog.png", Inches(0.4), Inches(2.0), Inches(6.1), Inches(4.6))
except: pass
try:
    add_image(s, SS/"04_board.png",   Inches(6.8), Inches(2.0), Inches(6.1), Inches(4.6))
except: pass

# Labels
add_text(s, "Backlog Görünümü",
         Inches(0.4), Inches(6.65), Inches(6.1), Inches(0.3),
         font_size=Pt(8.5), bold=True, color=BLUE_DK, align=PP_ALIGN.CENTER)
add_text(s, "Kanban Board",
         Inches(6.8), Inches(6.65), Inches(6.1), Inches(0.3),
         font_size=Pt(8.5), bold=True, color=BLUE_DK, align=PP_ALIGN.CENTER)

slide_footer(s, 4, TOTAL)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — QUALITY & TRACEABILITY
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_stripe(s)
section_title(s, "Kalite Yönetimi & Teslimat İzlenebilirliği")
add_text(s, "Test operasyonu ürünün çekirdeğinde — eklenti değil. Gereksinimden deployment'a otomatik zincir.",
         Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.35),
         font_size=Pt(10.5), color=BODY)

try:
    add_image(s, SS/"05_quality_catalog.png", Inches(0.4), Inches(2.0), Inches(6.1), Inches(4.6))
except: pass
try:
    add_image(s, SS/"06_traceability.png",    Inches(6.8), Inches(2.0), Inches(6.1), Inches(4.6))
except: pass

add_text(s, "Test Kataloğu & Kampanya Yönetimi",
         Inches(0.4), Inches(6.65), Inches(6.1), Inches(0.3),
         font_size=Pt(8.5), bold=True, color=BLUE_DK, align=PP_ALIGN.CENTER)
add_text(s, "Kalite İzlenebilirlik Matrisi",
         Inches(6.8), Inches(6.65), Inches(6.1), Inches(0.3),
         font_size=Pt(8.5), bold=True, color=BLUE_DK, align=PP_ALIGN.CENTER)

slide_footer(s, 5, TOTAL)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DASHBOARD
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_stripe(s)
section_title(s, "Yönetim Dashboard'u")
add_text(s, "Sprint metrikleri, kalite durumu ve teslimat görünürlüğü tek ekrandan — rapor toplamak gerekmez.",
         Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.35),
         font_size=Pt(10.5), color=BODY)

try:
    add_image(s, SS/"02_dashboard.png", Inches(1.2), Inches(2.0), Inches(10.9), Inches(5.0))
except: pass

slide_footer(s, 6, TOTAL)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 4 BUSINESS OUTCOMES
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_stripe(s)
section_title(s, "Dört Temel İş Kazanımı")

ow = Inches(6.1)
oh = Inches(2.4)
ogap = Inches(0.2)
ot = Inches(1.8)

outcome_card(s, Inches(0.4), ot, ow, oh,
    "Yönetim Görünürlüğü",
    "Sprint, kalite ve teslimat durumunu anlamak için araçlar arası veri toplama saatler alıyor.",
    "Dashboard'dan tek bakışta: aktif cycle durumu, açık defect sayısı, son deployment geçmişi.",
    "Karar gecikmesi ve yanlış beklenti riski azalır")

outcome_card(s, Inches(0.4)+ow+ogap, ot, ow, oh,
    "Kalite Operasyonu — Ürünün Merkezinde",
    "Test yönetimi ayrı araç, ayrı lisans. Gereksinim-test bağlantısı el emeğiyle kuruluyor.",
    "Test case, defect ve backlog aynı veri modelinde. Soru saniyeler içinde yanıtlanır.",
    "Kalite karar döngüsü kısalır; ek araç maliyeti ortadan kalkar")

ot2 = ot + oh + ogap
outcome_card(s, Inches(0.4), ot2, ow, oh,
    "Gereksinimden Teslimata Uçtan Uca Zincir",
    "'Bu değişiklik hangi ortama kadar gitti?' sorusu araçlar arası araştırma gerektiriyor.",
    "Commit ve deployment event'leri ilgili gereksinime otomatik bağlanır.",
    "Denetim ve uyum hazırlık süresi önemli ölçüde kısalır")

outcome_card(s, Inches(0.4)+ow+ogap, ot2, ow, oh,
    "Sürece Uyan Araç Modeli",
    "Sabit şablon tüm projelere aynı iş akışını dayatıyor. Ekip araca uyum sağlıyor.",
    "İş akışı, alan tanımları ve erişim kuralları yapılandırma dosyası ile tanımlanır.",
    "Farklı ekipler aynı platformda, kendi süreç standartlarında çalışır")

slide_footer(s, 7, TOTAL)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FINANCIAL CASE
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_stripe(s)
section_title(s, "Finansal Çerçeve")
add_text(s, "30–50 kişilik yazılım ekipleri için sektör araştırma raporlarından derlenen temsili maliyet karşılaştırması.",
         Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.35),
         font_size=Pt(10.5), color=BODY)

# Table header
th = Inches(1.95)
add_rect(s, Inches(0.5), th, Inches(12.3), Inches(0.42), fill=INK)
for text, lx, ww, align in [
    ("Maliyet Kalemi",         Inches(0.55), Inches(5.5), PP_ALIGN.LEFT),
    ("Mevcut Araç Zinciri",    Inches(6.1),  Inches(3.2), PP_ALIGN.CENTER),
    ("Pamera ile",             Inches(9.4),  Inches(3.3), PP_ALIGN.CENTER),
]:
    add_text(s, text, lx, th+Inches(0.07), ww, Inches(0.3),
             font_size=Pt(9), bold=True, color=WHITE, align=align)
# blue highlight on last col header
add_rect(s, Inches(9.35), th, Inches(3.5), Inches(0.42), fill=BLUE)
add_text(s, "Pamera ile", Inches(9.4), th+Inches(0.07), Inches(3.3), Inches(0.3),
         font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Backlog & iş takibi lisansı",    "$4.000–6.000 / yıl",   "Dahil"),
    ("Test yönetimi lisansı",           "$5.000–8.000 / yıl",   "Dahil"),
    ("Entegrasyon kurulum & bakım",     "$3.000–6.000 / yıl",   "Yok"),
    ("Denetim & raporlama hazırlığı",   "~$4.000 / yıl",        "Önemli ölçüde azalır"),
]

for i, (label, theirs, ours) in enumerate(rows):
    ry = th + Inches(0.42) + i * Inches(0.5)
    bg = BG_ALT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.5), ry, Inches(12.3), Inches(0.5), fill=bg)
    add_text(s, label,  Inches(0.6), ry+Inches(0.1), Inches(5.4), Inches(0.35),
             font_size=Pt(9.5), color=INK)
    add_text(s, theirs, Inches(6.1), ry+Inches(0.1), Inches(3.2), Inches(0.35),
             font_size=Pt(9.5), color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    add_rect(s, Inches(9.35), ry, Inches(3.5), Inches(0.5),
             fill=RGBColor(0xF0,0xFF,0xF4))
    add_text(s, ours,   Inches(9.4), ry+Inches(0.1), Inches(3.3), Inches(0.35),
             font_size=Pt(9.5), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Total row
ry = th + Inches(0.42) + len(rows) * Inches(0.5)
add_rect(s, Inches(0.5), ry, Inches(12.3), Inches(0.5), fill=BLUE_LT)
add_text(s, "Toplam tahmini yıllık maliyet",
         Inches(0.6), ry+Inches(0.1), Inches(5.4), Inches(0.35),
         font_size=Pt(9.5), bold=True, color=INK)
add_text(s, "$16.000–24.000",
         Inches(6.1), ry+Inches(0.1), Inches(3.2), Inches(0.35),
         font_size=Pt(9.5), color=MUTED, align=PP_ALIGN.CENTER, italic=True)
add_rect(s, Inches(9.35), ry, Inches(3.5), Inches(0.5),
         fill=RGBColor(0xD0,0xE6,0xFF))
add_text(s, "Tek kalem",
         Inches(9.4), ry+Inches(0.1), Inches(3.3), Inches(0.35),
         font_size=Pt(9.5), bold=True, color=BLUE_DK, align=PP_ALIGN.CENTER)

add_text(s, "Kaynak: Gartner Software Engineering Survey 2023 · Asana Anatomy of Work Report 2024 · IBM Systems Sciences Institute. "
            "Rakamlar sektör ortalamasına dayalı temsili değerlerdir.",
         Inches(0.5), H-Inches(0.85), Inches(12.3), Inches(0.3),
         font_size=Pt(7.5), color=MUTED, italic=True)

slide_footer(s, 8, TOTAL)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EVALUATION APPROACH + CONTACT
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_stripe(s)
section_title(s, "Değerlendirme Yaklaşımı")
add_text(s, "Kademeli ve düşük riskli süreç — her aşama, bir sonraki adım için bilgiye dayalı karar almayı sağlar.",
         Inches(0.5), Inches(1.45), Inches(8.0), Inches(0.35),
         font_size=Pt(10.5), color=BODY)

phases = [
    ("Aşama 1", "1–2 Hafta",
     "Keşif & Uyum Analizi",
     "Mevcut araç zinciri, ekip büyüklüğü ve öncelikli kullanım senaryolarının birlikte değerlendirilmesi. "
     "Pamera'nın hangi sorunlara yanıt verdiğinin netleştirilmesi."),
    ("Aşama 2", "4–6 Hafta",
     "Pilot Proje",
     "Tek bir projede, mevcut araçlarla paralel çalışma. Gerçek iş yüküyle platform testi; "
     "ekip geri bildirimi ve süreç uyarlaması. Zorunlu geçiş yoktur."),
    ("Aşama 3", "Planlanır",
     "Kurumsal Uygulama",
     "Pilot sonuçlarına dayalı kurumsal geçiş planı. Ekip onboarding, süreç standardizasyonu "
     "ve geçiş yol haritası birlikte hazırlanır."),
]

for i, (phase, dur, title, desc) in enumerate(phases):
    py = Inches(2.05) + i * Inches(1.45)
    # marker box
    add_rect(s, Inches(0.5), py, Inches(1.1), Inches(1.2), fill=BLUE_LT,
             line=RGBColor(0xB3,0xD4,0xFF), line_w=Pt(0.5))
    add_text(s, phase, Inches(0.5), py+Inches(0.08), Inches(1.1), Inches(0.3),
             font_size=Pt(7), bold=True, color=BLUE_DK, align=PP_ALIGN.CENTER)
    add_text(s, dur, Inches(0.5), py+Inches(0.38), Inches(1.1), Inches(0.3),
             font_size=Pt(8.5), bold=True, color=BLUE_DK, align=PP_ALIGN.CENTER)
    # content
    add_text(s, title, Inches(1.75), py+Inches(0.05), Inches(6.5), Inches(0.35),
             font_size=Pt(11), bold=True, color=INK)
    add_text(s, desc, Inches(1.75), py+Inches(0.38), Inches(6.5), Inches(0.8),
             font_size=Pt(9.5), color=BODY)
    # connector line
    if i < len(phases) - 1:
        add_rect(s, Inches(1.02), py+Inches(1.2), Inches(0.06), Inches(0.25), fill=BLUE_LT)

# Right panel — contact
add_rect(s, Inches(9.2), Inches(1.8), Inches(3.8), Inches(5.0),
         fill=INK)
add_rect(s, Inches(9.2), Inches(1.8), Inches(0.07), Inches(5.0), fill=BLUE)

add_text(s, "İletişim",
         Inches(9.4), Inches(2.0), Inches(3.4), Inches(0.4),
         font_size=Pt(11), bold=True, color=WHITE)

try:
    logo_path = DOCS / "provera_logo.png"
    if logo_path.exists():
        add_image(s, logo_path, Inches(9.4), Inches(2.5), Inches(2.0))
except: pass

add_text(s, "Provera\nTelekomünikasyon ve\nBilgi Teknolojileri",
         Inches(9.4), Inches(3.5), Inches(3.4), Inches(0.9),
         font_size=Pt(10), bold=True, color=WHITE)

add_text(s, "info@provera.com.tr",
         Inches(9.4), Inches(4.55), Inches(3.4), Inches(0.35),
         font_size=Pt(9.5), color=RGBColor(0xA0,0xC0,0xE8))
add_text(s, "provera.com.tr",
         Inches(9.4), Inches(4.9), Inches(3.4), Inches(0.35),
         font_size=Pt(9.5), color=RGBColor(0xA0,0xC0,0xE8))

add_text(s, "Gizli — Yalnızca\nAlıcı İçin",
         Inches(9.4), Inches(6.1), Inches(3.4), Inches(0.5),
         font_size=Pt(7.5), color=RGBColor(0x60,0x70,0x85), italic=True)

slide_footer(s, 9, TOTAL)

# ── Save ───────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Size:  {OUT.stat().st_size / 1024:.0f} KB")
